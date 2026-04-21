import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import requests
from PIL import Image

# --- 1. CONFIG & SECRETS ---
st.set_page_config(page_title="Plan B Media — LED Appendix Pro", layout="wide", page_icon="🎯")

st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;600;700;800&display=swap');
html, body, [class*="css"] { font-family: 'Inter', sans-serif; }

[data-testid="stSidebar"] {
    background: linear-gradient(180deg, #1B2A72 0%, #0D1B3E 100%) !important;
    border-right: 1px solid #00AEEF50;
}
[data-testid="stSidebar"] label,
[data-testid="stSidebar"] p,
[data-testid="stSidebar"] span,
[data-testid="stSidebar"] div { color: #FFFFFF !important; }
[data-testid="stSidebar"] hr { border-color: #00AEEF30 !important; }

[data-testid="stAppViewContainer"] {
    background: linear-gradient(160deg, #0D1B3E 0%, #0F2350 100%) !important;
}
[data-testid="stHeader"] { background: transparent !important; }

.planb-logo {
    padding: 24px 16px 18px;
    border-bottom: 2px solid #00AEEF;
    margin-bottom: 4px;
    text-align: center;
}
.planb-logo .brand {
    font-size: 28px; font-weight: 800; color: #fff;
    letter-spacing: -1px; line-height: 1;
}
.planb-logo .brand span { color: #00AEEF; }
.planb-logo .sub {
    font-size: 11px; color: #00AEEF;
    letter-spacing: 5px; text-transform: uppercase; margin-top: 3px;
}
.planb-logo .app-tag {
    display: inline-block; margin-top: 12px;
    background: #00AEEF20; border: 1px solid #00AEEF50;
    border-radius: 20px; padding: 3px 12px;
    font-size: 10px; color: #00AEEF; letter-spacing: 2px;
}

h1 {
    background: linear-gradient(90deg, #00AEEF 0%, #FFFFFF 60%);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
    background-clip: text; font-weight: 800 !important;
}
h2, h3 { color: #00AEEF !important; }

.stButton > button {
    border-radius: 8px !important; font-weight: 600 !important;
    border: 1px solid #00AEEF40 !important;
    transition: all 0.2s ease !important;
    color: #fff !important;
}
.stButton > button:hover {
    border-color: #00AEEF !important;
    box-shadow: 0 0 14px #00AEEF40 !important;
    transform: translateY(-1px) !important;
}
[data-testid="baseButton-primary"] {
    background: linear-gradient(135deg, #00AEEF, #0070A8) !important;
}

.stTextInput > div > div > input {
    background: #1B2A7250 !important;
    border: 1px solid #00AEEF40 !important;
    border-radius: 8px !important; color: white !important;
}
.stTextInput > div > div > input:focus {
    border-color: #00AEEF !important;
    box-shadow: 0 0 0 2px #00AEEF30 !important;
}

[data-testid="stMetric"] {
    background: #1B2A7260 !important;
    border: 1px solid #00AEEF30 !important;
    border-radius: 10px !important; padding: 10px !important;
}
[data-testid="stMetricValue"] { color: #00AEEF !important; font-weight: 800 !important; }

[data-testid="stVerticalBlockBorderWrapper"] {
    border: 1px solid #00AEEF20 !important;
    border-radius: 10px !important;
    background: #1B2A7215 !important;
    transition: border-color 0.2s, box-shadow 0.2s !important;
}
[data-testid="stVerticalBlockBorderWrapper"]:hover {
    border-color: #00AEEF70 !important;
    box-shadow: 0 0 12px #00AEEF20 !important;
}

[data-testid="stTabs"] [role="tab"] { font-weight: 600 !important; color: #ffffff80 !important; }
[data-testid="stTabs"] [role="tab"][aria-selected="true"] {
    color: #00AEEF !important; border-bottom-color: #00AEEF !important;
}

[data-testid="stAlert"] { border-radius: 10px !important; border-left-width: 4px !important; }

::-webkit-scrollbar { width: 5px; height: 5px; }
::-webkit-scrollbar-track { background: #0D1B3E; }
::-webkit-scrollbar-thumb { background: #00AEEF50; border-radius: 3px; }
::-webkit-scrollbar-thumb:hover { background: #00AEEF; }

hr { border-color: #00AEEF20 !important; }
</style>
""", unsafe_allow_html=True)

secrets = st.secrets if hasattr(st, 'secrets') else {}
API_KEY = secrets.get("GDRIVE_API_KEY", "")
FOLDER_ID = secrets.get("GDRIVE_FOLDER_ID", "")

# --- 2. SESSION STATE ---
if 'selected_images' not in st.session_state: st.session_state.selected_images = set()
if 'image_order' not in st.session_state: st.session_state.image_order = []
if 'drive_files' not in st.session_state: st.session_state.drive_files = []

# --- 3. FUNCTIONS ---
@st.cache_data(show_spinner="🚀 กำลังดึงข้อมูลจาก Drive...", ttl=3600)
def fetch_all_files(api_key, folder_id):
    if not api_key or not folder_id: return []
    all_files, page_token = [], None
    try:
        while True:
            params = {
                "q": f"'{folder_id}' in parents and mimeType contains 'image/' and trashed=false",
                "key": api_key,
                "fields": "nextPageToken,files(id,name,thumbnailLink)",
                "pageSize": 1000,
            }
            r = requests.get("https://www.googleapis.com/drive/v3/files", params=params, timeout=15)
            r.raise_for_status()
            data = r.json()
            all_files.extend(data.get("files", []))
            page_token = data.get("nextPageToken")
            if not page_token: break
        return sorted(all_files, key=lambda x: x["name"])
    except: return []

@st.cache_data(show_spinner=False)
def download_image(file_id, api_key):
    try:
        r = requests.get(f"https://www.googleapis.com/drive/v3/files/{file_id}",
                         params={"alt": "media", "key": api_key}, timeout=30)
        return r.content
    except: return None

def add_image_to_slide(slide, img_bytes, prs, is_cover=False):
    if not img_bytes: return
    try:
        if is_cover:
            slide.shapes.add_picture(io.BytesIO(img_bytes), 0, 0, width=prs.slide_width, height=prs.slide_height)
        else:
            img = Image.open(io.BytesIO(img_bytes))
            w, h = img.size
            ratio = min(prs.slide_width/w, prs.slide_height/h)
            nw, nh = int(w*ratio), int(h*ratio)
            slide.shapes.add_picture(io.BytesIO(img_bytes), (prs.slide_width-nw)//2, (prs.slide_height-nh)//2, width=nw, height=nh)
    except: pass

# --- 4. ⚡️ AUTO-LOAD ---
if API_KEY and FOLDER_ID and not st.session_state.drive_files:
    st.session_state.drive_files = fetch_all_files(API_KEY, FOLDER_ID)
    if st.session_state.drive_files: st.rerun()

# --- 5. LOGIC: MOVE & SORT ---
def move_left(idx):
    if idx > 0:
        st.session_state.image_order[idx], st.session_state.image_order[idx-1] = \
        st.session_state.image_order[idx-1], st.session_state.image_order[idx]

def move_right(idx):
    if idx < len(st.session_state.image_order) - 1:
        st.session_state.image_order[idx], st.session_state.image_order[idx+1] = \
        st.session_state.image_order[idx+1], st.session_state.image_order[idx]

def remove_item(fid):
    st.session_state.selected_images.discard(fid)
    st.session_state.image_order = [x for x in st.session_state.image_order if x != fid]

# --- 6. UI CONTENT ---
st.markdown("""
<div style="padding:8px 0 20px">
    <h1 style="margin:0;font-size:2.2rem">LED Appendix Pro</h1>
    <p style="color:#00AEEF90;margin:4px 0 0;font-size:.85rem;letter-spacing:2px;text-transform:uppercase">
        Plan B Media · Digital LED Proposal Builder
    </p>
</div>
""", unsafe_allow_html=True)
all_files = st.session_state.drive_files

if not all_files:
    st.error("❌ ไม่พบข้อมูล! โปรดเช็ก Secrets")
    if st.sidebar.button("🔄 ลองโหลดอีกครั้ง"):
        st.cache_data.clear()
        st.rerun()
    st.stop()

# ค้นหาปก COVER A / COVER B
def find_cover(key):
    for f in all_files:
        if key in f['name'].upper().replace(" ", ""): return f
    return None

cover_a, cover_b = find_cover("COVERA"), find_cover("COVERB")

with st.expander("📝 ตั้งค่าข้อความบนหน้าปก", expanded=False):
    c1, c2 = st.columns(2)
    p_title = c1.text_input("หัวข้อโปรเจกต์", "PROPOSAL FOR DIGITAL LED")
    p_sub = c2.text_input("ชื่อลูกค้า", "Presented to: Valued Customer")

tab1, tab2 = st.tabs(["📷 เลือกรูปภาพ", f"🖼️ จัดเรียงรูปภาพ ({len(st.session_state.selected_images)})"])

# --- TAB 1: เลือกรูป ---
with tab1:
    search = st.text_input("🔍 ค้นหารหัสภาพ", placeholder="พิมพ์ชื่อภาพเพื่อค้นหา...").upper()
    display = [f for f in all_files if (search and search in f["name"].upper()) or (f["id"] in st.session_state.selected_images)]
    
    if not display: st.info("💡 พิมพ์รหัสภาพเพื่อเริ่มค้นหา")
    else:
        cols = st.columns(4)
        for i, f in enumerate(display):
            with cols[i % 4]:
                with st.container(border=True):
                    st.image(f.get("thumbnailLink", "").replace("=s220", "=s400"), use_container_width=True)
                    fid, fname = f["id"], f["name"]
                    is_sel = fid in st.session_state.selected_images
                    if st.checkbox(fname[:15], key=f"sel_{fid}", value=is_sel):
                        if fid not in st.session_state.selected_images:
                            st.session_state.selected_images.add(fid)
                            st.session_state.image_order.append(fid)
                    else:
                        if fid in st.session_state.selected_images:
                            st.session_state.selected_images.discard(fid)
                            st.session_state.image_order = [x for x in st.session_state.image_order if x != fid]

# --- TAB 2: จัดเรียง (เวอร์ชันเรียงเป็นรูป) ---
with tab2:
    # กรองเฉพาะ ID ที่ยังถูกเลือกอยู่
    st.session_state.image_order = [fid for fid in st.session_state.image_order if fid in st.session_state.selected_images]
    order = st.session_state.image_order
    
    if not order:
        st.info("ยังไม่ได้เลือกรูปภาพ กรุณาเลือกจากแท็บแรกก่อนครับ")
    else:
        st.subheader("🖼️ จัดลำดับสไลด์ (กดลูกศรเพื่อสลับตำแหน่ง)")
        file_map = {f['id']: f for f in all_files}
        
        # แสดงรูปในรูปแบบ Grid และมีปุ่มเลื่อน ซ้าย/ขวา
        grid_cols = st.columns(3) # โชว์ทีละ 3 รูปต่อแถวเพื่อให้เห็นชัดบน iPad
        for i, fid in enumerate(order):
            f = file_map.get(fid)
            with grid_cols[i % 3]:
                with st.container(border=True):
                    st.image(f.get("thumbnailLink", "").replace("=s220", "=s400"), caption=f"ลำดับที่ {i+1}")
                    st.caption(f["name"])
                    # ปุ่มควบคุม
                    b1, b2, b3 = st.columns([1,1,1])
                    if b1.button("⬅️", key=f"L_{fid}_{i}", disabled=(i==0)):
                        move_left(i); st.rerun()
                    if b2.button("❌", key=f"RM_{fid}_{i}"):
                        remove_item(fid); st.rerun()
                    if b3.button("➡️", key=f"R_{fid}_{i}", disabled=(i==len(order)-1)):
                        move_right(i); st.rerun()

        st.divider()
        if st.button("🏗️ สร้างไฟล์ PowerPoint (.pptx)", type="primary", use_container_width=True):
            prs = Presentation()
            prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
            
            # 1. ปกหน้า
            slide_a = prs.slides.add_slide(prs.slide_layouts[6])
            if cover_a: add_image_to_slide(slide_a, download_image(cover_a['id'], API_KEY), prs, is_cover=True)
            tb = slide_a.shapes.add_textbox(0, prs.slide_height/2.5, prs.slide_width, Inches(3))
            p = tb.text_frame.paragraphs[0]
            p.text, p.alignment, p.font.size, p.font.bold = p_title, PP_ALIGN.CENTER, Pt(50), True
            if cover_a: p.font.color.rgb = RGBColor(255, 255, 255)
            p2 = tb.text_frame.add_paragraph()
            p2.text, p2.alignment, p2.font.size = p_sub, PP_ALIGN.CENTER, Pt(28)
            if cover_a: p2.font.color.rgb = RGBColor(255, 255, 255)

            # 2. เนื้อหา (เรียงตามที่กดสลับไว้)
            for fid in st.session_state.image_order:
                if (cover_a and fid == cover_a['id']) or (cover_b and fid == cover_b['id']): continue
                add_image_to_slide(prs.slides.add_slide(prs.slide_layouts[6]), download_image(fid, API_KEY), prs)

            # 3. ปกหลัง
            if cover_b:
                slide_b = prs.slides.add_slide(prs.slide_layouts[6])
                add_image_to_slide(slide_b, download_image(cover_b['id'], API_KEY), prs, is_cover=True)
            
            out = io.BytesIO()
            prs.save(out)
            st.success("🎉 Export สำเร็จ!")
            st.download_button("📥 ดาวน์โหลด Appendix", out.getvalue(), "Appendix.pptx", use_container_width=True)

# --- 7. SIDEBAR ---
with st.sidebar:
    st.markdown("""
    <div class="planb-logo">
        <div class="brand">Plan<span>·B</span></div>
        <div class="sub">media</div>
        <div class="app-tag">⚡ LED APPENDIX PRO</div>
    </div>
    """, unsafe_allow_html=True)
    st.subheader("📊 ข้อมูลปัจจุบัน")
    st.write(f"🖼️ ปกหน้า (COVER A): {'✅ พบแล้ว' if cover_a else '❌ ไม่พบ'}")
    st.write(f"🖼️ ปกหลัง (COVER B): {'✅ พบแล้ว' if cover_b else '❌ ไม่พบ'}")
    if st.button("🗑️ ล้าง Cache/เลือกใหม่"):
        st.cache_data.clear()
        st.session_state.selected_images = set()
        st.session_state.image_order = []
        st.rerun()
